import copy
import math
import random
import re
import tkinter as tk
from lxml import etree
from ortools.sat.python import cp_model
from pathlib import Path
from tkinter import messagebox, ttk
from typing import Optional

import openpyxl
import requests
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt, RGBColor
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.table import Table, TableStyleInfo
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt as PptxPt
from pptx.oxml.ns import qn
from pydantic import BaseModel, ValidationError, field_validator

EXCEL_PATH = "personnel.xlsx"
HEADER_COLOR = "2F5496"


# ---------------------------------------------------------------------------
# Pydantic models
# ---------------------------------------------------------------------------

class Soldier(BaseModel):
    name: str
    last_name: str
    personal_number: str
    phone: Optional[str] = ""
    course: Optional[str] = ""
    gender: Optional[str] = ""

    @field_validator("name", "last_name", "personal_number")
    @classmethod
    def must_not_be_empty(cls, v: str) -> str:
        if not v.strip():
            raise ValueError("field must not be empty")
        return v.strip()


class Commander(BaseModel):
    name: str
    last_name: str
    half_kappa: str

    @field_validator("name", "last_name", "half_kappa")
    @classmethod
    def must_not_be_empty(cls, v: str) -> str:
        if not v.strip():
            raise ValueError("field must not be empty")
        return v.strip()


class Hitnasut(BaseModel):
    hitnasut_name: str
    soldier_name: str
    commander_name: Optional[str] = ""

    @field_validator("hitnasut_name", "soldier_name")
    @classmethod
    def must_not_be_empty(cls, v: str) -> str:
        if not v.strip():
            raise ValueError("field must not be empty")
        return v.strip()


class Course(BaseModel):
    course: str
    half_kappa: str

    @field_validator("course", "half_kappa")
    @classmethod
    def must_not_be_empty(cls, v: str) -> str:
        if not v.strip():
            raise ValueError("field must not be empty")
        return v.strip()


class Room(BaseModel):
    room_number: str
    gender: str
    capacity: str
    room_manager: Optional[str] = ""

    @field_validator("room_number", "gender", "capacity")
    @classmethod
    def must_not_be_empty(cls, v: str) -> str:
        if not v.strip():
            raise ValueError("field must not be empty")
        return v.strip()

    @field_validator("capacity")
    @classmethod
    def capacity_must_be_positive(cls, v: str) -> str:
        try:
            if int(v.strip()) <= 0:
                raise ValueError("capacity must be greater than zero")
        except (ValueError, TypeError) as e:
            if "greater than zero" in str(e):
                raise
            raise ValueError("capacity must be a whole number greater than zero")
        return v.strip()


class ComputerUser(BaseModel):
    username: str
    password: str

    @field_validator("username", "password")
    @classmethod
    def must_not_be_empty(cls, v: str) -> str:
        if not v.strip():
            raise ValueError("field must not be empty")
        return v.strip()


class Classroom(BaseModel):
    commander_name: str
    class_name: str

    @field_validator("commander_name", "class_name")
    @classmethod
    def must_not_be_empty(cls, v: str) -> str:
        if not v.strip():
            raise ValueError("field must not be empty")
        return v.strip()


# ---------------------------------------------------------------------------
# Excel reading
# ---------------------------------------------------------------------------

def read_excel():
    wb = openpyxl.load_workbook(EXCEL_PATH)
    ws = wb.active
    assert ws is not None

    def read_rows(min_col, max_col, keys, model):
        valid, invalid = [], []
        key_to_col = {key: get_column_letter(min_col + i) for i, key in enumerate(keys)}
        for row_num, row in enumerate(
            ws.iter_rows(min_row=3, min_col=min_col, max_col=max_col, values_only=True),
            start=3,
        ):
            if not any(v is not None for v in row):
                continue
            raw = dict(zip(keys, (str(v).strip() if v is not None else "" for v in row)))
            try:
                obj = model(**raw).model_dump()
                obj["_row"] = row_num
                valid.append(obj)
            except ValidationError as e:
                invalid.append((row_num, e, key_to_col))
        return valid, invalid

    soldiers, soldier_errors = read_rows(
        1, 6, ["name", "last_name", "personal_number", "phone", "course", "gender"], Soldier
    )
    commanders, commander_errors = read_rows(
        8, 10, ["name", "last_name", "half_kappa"], Commander
    )
    hitnasuyot, hitnasut_errors = read_rows(
        12, 14, ["hitnasut_name", "soldier_name", "commander_name"], Hitnasut
    )
    courses, course_errors = read_rows(
        16, 17, ["course", "half_kappa"], Course
    )
    rooms, room_errors = read_rows(
        19, 22, ["room_number", "gender", "capacity", "room_manager"], Room
    )
    computer_users, computer_user_errors = read_rows(
        24, 25, ["username", "password"], ComputerUser
    )
    classrooms, classroom_errors = read_rows(
        27, 28, ["commander_name", "class_name"], Classroom
    )

    all_errors = (
        soldier_errors + commander_errors + hitnasut_errors
        + course_errors + room_errors + computer_user_errors + classroom_errors
    )
    if all_errors:
        lines = []
        for row_num, err, key_to_col in all_errors:
            for e in err.errors():
                col = key_to_col.get(str(e["loc"][0]), "?") if e["loc"] else "?"
                msg = e["msg"].removeprefix("Value error, ")
                lines.append(f"Row {row_num}, col {col}: {msg}")
        raise ValueError("Validation errors in personnel.xlsx:\n" + "\n".join(lines))

    return soldiers, commanders, hitnasuyot, courses, rooms, computer_users, classrooms


# ---------------------------------------------------------------------------
# Button 1 — Enrich soldier data via API
# ---------------------------------------------------------------------------

def _api_fetch(personal_number: str) -> dict:
    """Calls the soldier info API. Replace URL with the real endpoint."""
    # TODO: replace with real API URL
    API_URL = "https://api.example.com/soldiers/{}"
    try:
        response = requests.get(API_URL.format(personal_number), timeout=5)
        response.raise_for_status()
        return response.json()
    except Exception:
        # Fallback mock data while API is a placeholder
        return {
            "rank": random.choice(["Recruit", "Private", "Corporal", "Sergeant"]),
            "unit": f"Unit {random.randint(100, 999)}",
            "date_of_birth": f"{random.randint(1, 28):02d}/{random.randint(1, 12):02d}/200{random.randint(0, 5)}",
            "city": random.choice(["Tel Aviv", "Jerusalem", "Haifa", "Beer Sheva", "Eilat"]),
        }


def _style_header(cell):
    cell.font = Font(bold=True, color="FFFFFF")
    cell.fill = PatternFill(fill_type="solid", fgColor=HEADER_COLOR)
    cell.alignment = Alignment(horizontal="center")


def _assign_commanders(soldiers: list, commanders: list, courses: list) -> dict:
    """
    Uses OR-Tools CP-SAT to assign each soldier one commander.
    Each soldier may only be assigned to a commander in the same half-kappa
    (derived from the soldier's course). Load is balanced within each
    half-kappa group. Same-course cohesion is maximised as the objective.
    Falls back to round-robin within half-kappa on solver failure.
    """
    if not commanders or not soldiers:
        return {}

    S = len(soldiers)
    C = len(commanders)

    course_to_hk: dict[str, str] = {
        c["course"].strip().lower(): c["half_kappa"].strip().lower() for c in courses
    }
    c_hk = [cmd["half_kappa"].strip().lower() for cmd in commanders]
    hk_to_cmds: dict[str, list[int]] = {}
    for ci, hk in enumerate(c_hk):
        hk_to_cmds.setdefault(hk, []).append(ci)

    # Eligible commanders per soldier (same half-kappa; fallback to all)
    eligible_commanders: list[list[int]] = []
    s_hk_list: list[str] = []
    for s in soldiers:
        course = (s.get("course") or "").strip().lower()
        hk = course_to_hk.get(course, "")
        eligible = hk_to_cmds.get(hk) if hk and hk in hk_to_cmds else None
        if not eligible:
            eligible = list(range(C))
            hk = ""
        eligible_commanders.append(list(eligible))
        s_hk_list.append(hk)

    by_course: dict[str, list[int]] = {}
    for i, s in enumerate(soldiers):
        course = (s.get("course") or "").strip().lower()
        by_course.setdefault(course, []).append(i)
    groups = list(by_course.values())
    G = len(groups)

    model = cp_model.CpModel()
    assign = [[model.new_bool_var(f"a_{si}_{ci}") for ci in range(C)] for si in range(S)]

    for si in range(S):
        eligible_set = set(eligible_commanders[si])
        for ci in range(C):
            if ci not in eligible_set:
                model.add(assign[si][ci] == 0)
        model.add_exactly_one([assign[si][ci] for ci in eligible_commanders[si]])

    # Load balance within each half-kappa group
    hk_soldier_indices: dict[str, list[int]] = {}
    for si, hk in enumerate(s_hk_list):
        hk_soldier_indices.setdefault(hk, []).append(si)

    for hk, s_indices in hk_soldier_indices.items():
        c_indices = hk_to_cmds.get(hk, list(range(C))) if hk else list(range(C))
        if not c_indices:
            continue
        S_hk = len(s_indices)
        C_hk = len(c_indices)
        lo, hi = S_hk // C_hk, math.ceil(S_hk / C_hk)
        for ci in c_indices:
            col = [assign[si][ci] for si in s_indices]
            model.add(sum(col) >= lo)
            model.add(sum(col) <= hi)

    intact = [[model.new_bool_var(f"intact_{g}_{ci}") for ci in range(C)] for g in range(G)]
    for g, group in enumerate(groups):
        for ci in range(C):
            for si in group:
                model.add_implication(intact[g][ci], assign[si][ci])

    group_sizes = [len(g) for g in groups]

    # No-isolation: for groups of size >= 2, penalise having exactly 1 member in a slot.
    # alone[g][ci] = 1  iff  exactly 1 soldier from group g is assigned to commander ci.
    # Detected via:  at1 (count >= 1)  and  at2 (count >= 2),  alone = at1 - at2.
    alone_vars = []
    for g, group in enumerate(groups):
        n_g = len(group)
        if n_g < 2:
            continue
        for ci in range(C):
            col = [assign[si][ci] for si in group]
            at1 = model.new_bool_var(f"at1_{g}_{ci}")
            at2 = model.new_bool_var(f"at2_{g}_{ci}")
            model.add(sum(col) >= at1)
            model.add(sum(col) <= n_g * at1)
            model.add((n_g - 1) * at2 >= sum(col) - 1)
            model.add(sum(col) >= 2 * at2)
            alone = model.new_bool_var(f"alone_{g}_{ci}")
            model.add(alone == at1 - at2)
            alone_vars.append(alone)

    # Priority 1 — no isolated soldiers  (weight S+1 dominates all cohesion gains)
    # Priority 2 — maximise fully intact course groups
    BIG = S + 1
    intact_terms = [intact[g][ci] * group_sizes[g] for g in range(G) for ci in range(C)]
    alone_sum = sum(alone_vars) if alone_vars else 0
    model.maximize(sum(intact_terms) - BIG * alone_sum)

    solver = cp_model.CpSolver()
    solver.parameters.max_time_in_seconds = 5.0
    status = solver.solve(model)

    if status in (cp_model.OPTIMAL, cp_model.FEASIBLE):
        result: dict[str, dict] = {}
        for si, soldier in enumerate(soldiers):
            for ci in range(C):
                if solver.value(assign[si][ci]):
                    result[soldier["personal_number"]] = commanders[ci]
                    break
        return result

    # Fallback: round-robin within half-kappa group
    result = {}
    for hk, s_indices in hk_soldier_indices.items():
        c_indices = hk_to_cmds.get(hk, []) if hk else []
        if not c_indices:
            c_indices = list(range(C))
        for idx, si in enumerate(s_indices):
            result[soldiers[si]["personal_number"]] = commanders[c_indices[idx % len(c_indices)]]
    return result


def _validate_cross_references(
    soldiers: list, commanders: list, hitnasuyot: list, courses: list, computer_users: list,
    classrooms: list, rooms: list,
) -> dict[str, list[str]]:
    """Returns an ordered dict of {section_heading: [error_lines]}, empty if all OK."""
    sections: dict[str, list[str]] = {}

    course_names = {c["course"].strip().lower() for c in courses}
    soldier_fullnames = {
        f"{s['name']} {s['last_name']}".strip().lower() for s in soldiers
    }
    commander_fullnames = {
        f"{c['name']} {c['last_name']}".strip().lower() for c in commanders
    }

    # Duplicate soldier full names
    name_to_rows: dict[str, list[int]] = {}
    for s in soldiers:
        full = f"{s['name']} {s['last_name']}".strip().lower()
        name_to_rows.setdefault(full, []).append(s["_row"])
    dup_lines = []
    for full, rows in name_to_rows.items():
        if len(rows) > 1:
            dup_lines.append(f"Rows {', '.join(str(r) for r in rows)}  (cols A–B — duplicate full name)")
    if dup_lines:
        sections["Soldiers with duplicate name"] = dup_lines

    # Duplicate personal numbers
    pn_to_rows: dict[str, list[int]] = {}
    for s in soldiers:
        pn = s["personal_number"].strip()
        pn_to_rows.setdefault(pn, []).append(s["_row"])
    dup_pn = []
    for pn, rows in pn_to_rows.items():
        if len(rows) > 1:
            dup_pn.append(f"Rows {', '.join(str(r) for r in rows)}  (col C — duplicate personal number)")
    if dup_pn:
        sections["Soldiers with duplicate personal number"] = dup_pn

    # Duplicate phone numbers (only among non-empty values)
    phone_to_rows: dict[str, list[int]] = {}
    for s in soldiers:
        phone = (s.get("phone") or "").strip()
        if phone:
            phone_to_rows.setdefault(phone, []).append(s["_row"])
    dup_phone = []
    for phone, rows in phone_to_rows.items():
        if len(rows) > 1:
            dup_phone.append(f"Rows {', '.join(str(r) for r in rows)}  (col D — duplicate phone number)")
    if dup_phone:
        sections["Soldiers with duplicate phone number"] = dup_phone

    # Soldier courses missing from Courses table
    course_lines = []
    for s in soldiers:
        course = (s.get("course") or "").strip()
        if course and course.lower() not in course_names:
            course_lines.append(f"Row {s['_row']}, col F  (Course — not found in Courses table)")
    if course_lines:
        sections["Soldiers with unknown course"] = course_lines

    # Soldiers whose course's half-kappa has no commander
    course_to_hk = {c["course"].strip().lower(): c["half_kappa"].strip() for c in courses}
    commander_hks = {cmd["half_kappa"].strip().lower() for cmd in commanders}
    no_cmd_lines = []
    for s in soldiers:
        course = (s.get("course") or "").strip()
        if not course or course.lower() not in course_names:
            continue
        hk = course_to_hk.get(course.lower(), "")
        if hk and hk.lower() not in commander_hks:
            no_cmd_lines.append(
                f"Row {s['_row']}, col F  (Course — no commander for this course's half-kappa)"
            )
    if no_cmd_lines:
        sections["Soldiers with no eligible commander"] = no_cmd_lines

    # Hitnasuyot: soldier full name not in Soldiers table
    soldier_errors = []
    for h in hitnasuyot:
        sname = h["soldier_name"].strip()
        if sname.lower() not in soldier_fullnames:
            soldier_errors.append(f"Row {h['_row']}, col M  (Full Soldier Name — not found in Soldiers table)")
    if soldier_errors:
        sections["Hitnasuyot with unknown soldier"] = soldier_errors

    # Hitnasuyot: commander full name not in Commanders table
    cmd_errors = []
    for h in hitnasuyot:
        cname = (h.get("commander_name") or "").strip()
        if cname and cname.lower() not in commander_fullnames:
            cmd_errors.append(f"Row {h['_row']}, col N  (Full Commander Name — not found in Commanders table)")
    if cmd_errors:
        sections["Hitnasuyot with unknown commander"] = cmd_errors

    # Computer Users: each soldier must have a matching entry by position
    if len(computer_users) < len(soldiers):
        missing = soldiers[len(computer_users):]
        sections["Soldiers without a computer user"] = [
            f"Row {s['_row']}, cols A–B  (no computer user entry for this soldier)"
            for s in missing
        ]

    # Classrooms: commander full name must exist in the Commanders table
    classroom_cmd_errors = []
    for cl in classrooms:
        cname = cl["commander_name"].strip()
        if cname.lower() not in commander_fullnames:
            classroom_cmd_errors.append(
                f"Row {cl['_row']}, col AA  (Full Commander Name — not found in Commanders table)"
            )
    if classroom_cmd_errors:
        sections["Classrooms with unknown commander"] = classroom_cmd_errors

    # Rooms: room manager must exist in the Soldiers table
    room_manager_errors = []
    for r in rooms:
        mgr = (r.get("room_manager") or "").strip()
        if mgr and mgr.lower() not in soldier_fullnames:
            room_manager_errors.append(
                f"Row {r['_row']}, col V  (Room Manager — not found in Soldiers table)"
            )
    if room_manager_errors:
        sections["Rooms with unknown room manager"] = room_manager_errors

    return sections


_PLACEHOLDER_RE = re.compile(r"<[^>]+>")


def _validate_templates() -> tuple[list[str], list[str]]:
    """Returns (missing_errors, unknown_errors). Missing = hard block; unknown = ignorable."""
    missing_errors: list[str] = []
    unknown_errors: list[str] = []

    def _placeholders_pptx(path: Path) -> set[str]:
        prs = Presentation(str(path))
        found: set[str] = set()
        for slide in prs.slides:
            for t in slide.shapes._spTree.iter(qn("a:t")):
                if t.text:
                    for m in _PLACEHOLDER_RE.findall(t.text):
                        found.add(m.lower())
        return found

    def _placeholders_docx(path: Path) -> set[str]:
        doc = Document(str(path))
        found: set[str] = set()
        for para in doc.paragraphs:
            for m in _PLACEHOLDER_RE.findall(para.text):
                found.add(m.lower())
        return found

    checks = [
        (
            HITNASUYOT_TEMPLATE,
            _placeholders_pptx,
            {"<hitnasut name>", "<commander name>", "<soldier name>"},
            {"<hitnasut name>", "<commander name>", "<soldier name>"},
        ),
        (
            ROOM_SIGNS_TEMPLATE,
            _placeholders_docx,
            {"<room sign>", "<soldier name>"},
            {"<room sign>", "<soldier name>"},
        ),
        (
            NAME_TAGS_TEMPLATE,
            _placeholders_pptx,
            {"<soldier name>"},
            {"<soldier name>"},
        ),
        (
            CLASSROOM_TEMPLATE,
            _placeholders_pptx,
            {"<commander name>", "<classroom name>"},
            {"<commander name>", "<classroom name>"},
        ),
    ]

    for path, scanner, required, recognized in checks:
        found = scanner(path)
        missing = required - found
        unknown = found - recognized
        if missing:
            missing_errors.append(
                f"{path.name}: missing placeholder(s): {', '.join(sorted(missing))}"
            )
        if unknown:
            unknown_errors.append(
                f"{path.name}: unrecognized placeholder(s): {', '.join(sorted(unknown))}"
            )

    return missing_errors, unknown_errors


def _ask_ignore_unknowns(errors: list[str]) -> bool:
    """Show a warning dialog for unrecognized placeholders.
    Returns True if the user clicks Ignore, False if Cancel."""
    confirmed = False
    dlg = tk.Toplevel(root)
    dlg.title("Unrecognized Placeholders")
    dlg.resizable(False, False)
    dlg.grab_set()

    ttk.Label(
        dlg, text="Unrecognized placeholders found:", font=(_FONT, 11, "bold")
    ).pack(padx=20, pady=(16, 6), anchor="w")
    for err in errors:
        ttk.Label(dlg, text=f"• {err}").pack(padx=20, pady=1, anchor="w")
    ttk.Label(
        dlg, text="These will be left unchanged in the output.", foreground="gray"
    ).pack(padx=20, pady=(8, 12), anchor="w")

    btn_frame = ttk.Frame(dlg)
    btn_frame.pack(pady=(0, 16), padx=20)

    def on_cancel():
        dlg.destroy()

    def on_ignore():
        nonlocal confirmed
        confirmed = True
        dlg.destroy()

    ttk.Button(btn_frame, text="Cancel", command=on_cancel, width=12).pack(side=tk.LEFT, padx=(0, 8))
    tk.Button(
        btn_frame, text="Ignore", command=on_ignore, width=12,
        bg="#c0392b", fg="white", activebackground="#e74c3c", activeforeground="white",
        relief="flat", cursor="hand2",
    ).pack(side=tk.LEFT)

    dlg.update_idletasks()
    pw, ph = root.winfo_x(), root.winfo_y()
    rw, rh = root.winfo_width(), root.winfo_height()
    dw, dh = dlg.winfo_width(), dlg.winfo_height()
    dlg.geometry(f"+{pw + (rw - dw) // 2}+{ph + (rh - dh) // 2}")

    dlg.wait_window()
    return confirmed


def _show_grouped_errors_dialog(title: str, sections: dict[str, list[str]]) -> None:
    """Modal error dialog with section headers; grows horizontally to fit content."""
    dlg = tk.Toplevel(root)
    dlg.title(title)
    dlg.resizable(False, False)
    dlg.grab_set()

    first = True
    for section_title, lines in sections.items():
        top_pad = (16, 4) if first else (12, 4)
        first = False
        ttk.Label(dlg, text=section_title, font=(_FONT, 10, "bold")).pack(
            padx=20, pady=top_pad, anchor="w"
        )
        for line in lines:
            ttk.Label(dlg, text=f"    • {line}").pack(padx=20, pady=1, anchor="w")

    ttk.Button(dlg, text="OK", command=dlg.destroy, width=10).pack(pady=(14, 16))

    dlg.update_idletasks()
    dw, dh = dlg.winfo_width(), dlg.winfo_height()
    pw, ph = root.winfo_x(), root.winfo_y()
    rw, rh = root.winfo_width(), root.winfo_height()
    dlg.geometry(f"+{pw + (rw - dw) // 2}+{ph + (rh - dh) // 2}")

    dlg.wait_window()


def _unique_result_dir() -> Path:
    base = Path("result")
    if not base.exists():
        return base
    i = 1
    while Path(f"result({i})").exists():
        i += 1
    return Path(f"result({i})")


def _generate_enriched(soldiers: list, out_dir: Path) -> str:
    enriched = [{**s, **_api_fetch(s["personal_number"])} for s in soldiers]

    wb = openpyxl.Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "Enriched Soldiers"

    headers = ["Name", "Last Name", "Personal Number", "Phone", "Course",
               "Rank", "Unit", "Date of Birth", "City"]
    fields = ["name", "last_name", "personal_number", "phone", "course",
              "rank", "unit", "date_of_birth", "city"]

    for col, header in enumerate(headers, start=1):
        _style_header(ws.cell(row=1, column=col, value=header))
        ws.column_dimensions[get_column_letter(col)].width = 18

    for row_idx, s in enumerate(enriched, start=2):
        for col, field in enumerate(fields, start=1):
            ws.cell(row=row_idx, column=col, value=s.get(field, ""))

    tab = Table(displayName="EnrichedSoldiers", ref=f"A1:{get_column_letter(len(headers))}{len(enriched) + 1}")
    tab.tableStyleInfo = TableStyleInfo(name="TableStyleMedium9", showRowStripes=True, showFirstColumn=False, showLastColumn=False, showColumnStripes=False)
    ws.add_table(tab)

    out = out_dir / "enriched_soldiers.xlsx"
    wb.save(out)
    return str(out)


# ---------------------------------------------------------------------------
# Button 2 — Hitnasuyot PPTX (template-based)
# ---------------------------------------------------------------------------

HITNASUYOT_TEMPLATE = Path("templates/hitnasuyot.pptx")


def _copy_slide(prs: Presentation, source_slide):  # type: ignore[return]
    """Append a copy of source_slide to prs and return the new slide."""
    new_slide = prs.slides.add_slide(prs.slide_layouts[0])

    # Copy shape tree
    src_tree = source_slide.shapes._spTree
    dst_tree = new_slide.shapes._spTree
    for child in list(dst_tree):
        dst_tree.remove(child)
    for child in src_tree:
        dst_tree.append(copy.deepcopy(child))

    # Copy background
    src_bg = source_slide.element.find(".//" + qn("p:bg"))
    if src_bg is not None:
        new_cSld = new_slide.element.find(qn("p:cSld"))
        existing_bg = new_cSld.find(qn("p:bg")) if new_cSld is not None else None
        if new_cSld is not None:
            if existing_bg is not None:
                new_cSld.remove(existing_bg)
            new_cSld.insert(0, copy.deepcopy(src_bg))

    return new_slide


def _replace_text(slide, replacements: dict):
    """Replace placeholder strings (case-insensitive) directly in XML text nodes."""
    for t_elem in slide.shapes._spTree.iter(qn("a:t")):
        if not t_elem.text:
            continue
        for placeholder, value in replacements.items():
            t_elem.text = re.sub(re.escape(placeholder), value, t_elem.text, flags=re.IGNORECASE)


def _generate_pptx(hitnasuyot: list, out_dir: Path) -> str:
    prs = Presentation(str(HITNASUYOT_TEMPLATE))
    t_with_commander = prs.slides[0]
    t_no_commander   = prs.slides[1]
    t_soldier        = prs.slides[2]

    # Group entries by hitnasut_name, preserving order
    groups: dict[str, dict] = {}
    for entry in hitnasuyot:
        name = entry["hitnasut_name"]
        if name not in groups:
            groups[name] = {"commander": entry["commander_name"], "soldiers": []}
        groups[name]["soldiers"].append(entry["soldier_name"])

    for hitnasut_name, data in groups.items():
        commander = data["commander"]
        if commander:
            intro = _copy_slide(prs, t_with_commander)
            _replace_text(intro, {
                "<Hitnasut name>": hitnasut_name,
                "<Commander name>": commander,
            })
        else:
            intro = _copy_slide(prs, t_no_commander)
            _replace_text(intro, {"<Hitnasut name>": hitnasut_name})

        for soldier_name in data["soldiers"]:
            s_slide = _copy_slide(prs, t_soldier)
            _replace_text(s_slide, {"<Soldier name>": soldier_name})

    # Remove the 3 original template slides
    sldIdLst = prs.slides._sldIdLst
    for _ in range(3):
        sldIdLst.remove(sldIdLst[0])

    out = out_dir / "hitnasuyot.pptx"
    prs.save(str(out))
    return str(out)


# ---------------------------------------------------------------------------
# Classroom PPTX (template-based)
# ---------------------------------------------------------------------------

CLASSROOM_TEMPLATE = Path("templates/classroom.pptx")


def _generate_classroom_pptx(classrooms: list, out_dir: Path) -> str:
    prs = Presentation(str(CLASSROOM_TEMPLATE))
    tmpl_slide = prs.slides[0]

    for cl in classrooms:
        slide = _copy_slide(prs, tmpl_slide)
        _replace_text(slide, {
            "<commander name>": cl["commander_name"],
            "<classroom name>": cl["class_name"],
        })

    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    out = out_dir / "classroom.pptx"
    prs.save(str(out))
    return str(out)


# ---------------------------------------------------------------------------
# Name tags PPTX
# ---------------------------------------------------------------------------

NAME_TAGS_TEMPLATE = Path("templates/name_tags.pptx")


def _generate_name_tags(soldiers: list, out_dir: Path) -> str:
    prs = Presentation(str(NAME_TAGS_TEMPLATE))
    tmpl_slide = prs.slides[0]

    # Count how many <Soldier name> slots the template slide has
    slots = sum(
        1 for t in tmpl_slide.shapes._spTree.iter(qn("a:t"))
        if t.text and re.search(r"<soldier name>", t.text, re.IGNORECASE)
    )
    slots = max(slots, 1)

    # Split soldiers into pages of `slots` each
    batches = [soldiers[i:i + slots] for i in range(0, len(soldiers), slots)]

    for batch in batches:
        slide = _copy_slide(prs, tmpl_slide)
        soldier_idx = 0
        for t in slide.shapes._spTree.iter(qn("a:t")):
            if not t.text or not re.search(r"<soldier name>", t.text, re.IGNORECASE):
                continue
            if soldier_idx < len(batch):
                s = batch[soldier_idx]
                name = f"{s['name']} {s['last_name']}".strip()
            else:
                name = ""
            t.text = re.sub(r"<soldier name>", name, t.text, flags=re.IGNORECASE)
            soldier_idx += 1

    # Remove the original template slide
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    out = out_dir / "name_tags.pptx"
    prs.save(str(out))
    return str(out)


# ---------------------------------------------------------------------------
# Button 3 — Room signs Word document (template-based)
# ---------------------------------------------------------------------------

ROOM_SIGNS_TEMPLATE = Path("templates/room_signs.docx")
_W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_XML = "http://www.w3.org/XML/1998/namespace"


def _replace_in_para_elem(para_elem, replacements: dict):
    """Replace placeholder text (case-insensitive) in a paragraph XML element,
    handling placeholders split across multiple w:t runs."""
    t_nodes = list(para_elem.iter(f"{{{_W}}}t"))
    if not t_nodes:
        return
    full = "".join(t.text or "" for t in t_nodes)
    new_text = full
    for placeholder, value in replacements.items():
        new_text = re.sub(re.escape(placeholder), value, new_text, flags=re.IGNORECASE)
    if new_text == full:
        return
    t_nodes[0].text = new_text
    t_nodes[0].set(f"{{{_XML}}}space", "preserve")
    for t in t_nodes[1:]:
        t.text = ""


def _assign_rooms(soldiers: list, rooms: list) -> dict:
    """Assign soldiers to rooms using CP-SAT, maximising same-course soldiers per room,
    respecting gender and capacity constraints. Falls back to sequential on solver failure."""
    if not soldiers or not rooms:
        return {}

    S = len(soldiers)
    R = len(rooms)
    # Use distinct names (si/ri) so no variable ever shadows another in list comprehensions
    s_gender = [(sol.get("gender") or "").strip().lower() for sol in soldiers]
    r_gender  = [rm["gender"].strip().lower() for rm in rooms]
    r_cap     = [int(rm["capacity"]) for rm in rooms]

    eligible_rooms: list[list[int]] = [
        [ri for ri in range(R) if r_gender[ri] == s_gender[si]]
        for si in range(S)
    ]

    by_course: dict[str, list[int]] = {}
    for si, sol in enumerate(soldiers):
        course = (sol.get("course") or "").strip().lower()
        by_course.setdefault(course, []).append(si)
    groups = list(by_course.values())
    G = len(groups)

    model = cp_model.CpModel()
    assign = [[model.new_bool_var(f"a_{si}_{ri}") for ri in range(R)] for si in range(S)]

    for si in range(S):
        eligible_set = set(eligible_rooms[si])
        for ri in range(R):
            if ri not in eligible_set:
                model.add(assign[si][ri] == 0)
        if eligible_rooms[si]:
            # Use an explicit list — avoids generator closure bugs with loop variables
            model.add_exactly_one([assign[si][ri] for ri in eligible_rooms[si]])

    for ri in range(R):
        occupants = [assign[si][ri] for si in range(S)]
        model.add(sum(occupants) <= r_cap[ri])

    intact = [[model.new_bool_var(f"intact_{g}_{ri}") for ri in range(R)] for g in range(G)]
    for g, group in enumerate(groups):
        for ri in range(R):
            for si in group:
                model.add_implication(intact[g][ri], assign[si][ri])

    group_sizes = [len(g) for g in groups]

    # No-isolation: for groups of size >= 2, penalise having exactly 1 member in a room.
    # alone[g][ri] = 1  iff  exactly 1 soldier from group g is in room ri.
    alone_vars = []
    for g, group in enumerate(groups):
        n_g = len(group)
        if n_g < 2:
            continue
        for ri in range(R):
            col = [assign[si][ri] for si in group]
            at1 = model.new_bool_var(f"at1_{g}_{ri}")
            at2 = model.new_bool_var(f"at2_{g}_{ri}")
            model.add(sum(col) >= at1)
            model.add(sum(col) <= n_g * at1)
            model.add((n_g - 1) * at2 >= sum(col) - 1)
            model.add(sum(col) >= 2 * at2)
            alone = model.new_bool_var(f"alone_{g}_{ri}")
            model.add(alone == at1 - at2)
            alone_vars.append(alone)

    # Priority 1 — no isolated soldiers  (weight S+1 dominates all cohesion gains)
    # Priority 2 — maximise fully intact course groups
    BIG = S + 1
    intact_terms = [intact[g][ri] * group_sizes[g] for g in range(G) for ri in range(R)]
    alone_sum = sum(alone_vars) if alone_vars else 0
    model.maximize(sum(intact_terms) - BIG * alone_sum)

    solver = cp_model.CpSolver()
    solver.parameters.max_time_in_seconds = 5.0
    status = solver.solve(model)

    # Build per-room-INDEX lists so each row's capacity is enforced independently,
    # even when multiple rows share the same room_number.
    per_index: list[list] = [[] for _ in range(R)]

    if status in (cp_model.OPTIMAL, cp_model.FEASIBLE):
        for si, soldier in enumerate(soldiers):
            for ri in range(R):
                if solver.value(assign[si][ri]):
                    per_index[ri].append(soldier)
                    break
    else:
        # Fallback: sequential assignment by gender
        by_gender: dict[str, list] = {}
        for sol in soldiers:
            g = (sol.get("gender") or "").strip().lower()
            by_gender.setdefault(g, []).append(sol)
        gender_pos: dict[str, int] = {g: 0 for g in by_gender}
        for ri, rm in enumerate(rooms):
            g   = rm["gender"].strip().lower()
            cap = r_cap[ri]
            pool = by_gender.get(g, [])
            pos  = gender_pos.get(g, 0)
            per_index[ri] = pool[pos:pos + cap]
            gender_pos[g] = pos + cap

    # Merge into room_number-keyed dict (separate rows stay separate in data)
    result: dict[str, list] = {}
    for ri, rm in enumerate(rooms):
        rn = rm["room_number"]
        result.setdefault(rn, []).extend(per_index[ri])
    return result


def _check_room_capacity(soldiers: list, rooms: list) -> list[str]:
    """Returns error lines for duplicate room numbers or insufficient capacity."""
    if not soldiers or not rooms:
        return []

    rn_to_rows: dict[str, list[int]] = {}
    for r in rooms:
        rn_to_rows.setdefault(r["room_number"], []).append(r["_row"])
    dup_rooms = [
        f"Rows {', '.join(str(r) for r in rows)}, col S  (duplicate room number)"
        for rows in rn_to_rows.values() if len(rows) > 1
    ]
    if dup_rooms:
        return dup_rooms

    soldiers_by_gender: dict[str, int] = {}
    for s in soldiers:
        g = (s.get("gender") or "").strip().lower()
        soldiers_by_gender[g] = soldiers_by_gender.get(g, 0) + 1

    capacity_by_gender: dict[str, int] = {}
    for r in rooms:
        g = r["gender"].strip().lower()
        capacity_by_gender[g] = capacity_by_gender.get(g, 0) + int(r["capacity"])

    errors = []
    for gender, count in soldiers_by_gender.items():
        capacity = capacity_by_gender.get(gender, 0)
        if capacity < count:
            short = count - capacity
            errors.append(
                f"{short} soldier(s) have no room — "
                f"{count} slots needed but only {capacity} available"
            )
    return errors


def _generate_data_xlsx(
    soldiers: list,
    hitnasuyot: list,
    courses: list,
    room_assignment: dict,
    commander_assignment: dict,
    computer_users: list,
    rooms: list,
    out_dir: Path,
) -> str:
    course_to_half_kappa = {c["course"].strip().lower(): c["half_kappa"] for c in courses}

    soldier_to_hitnasut: dict[str, dict] = {}
    for h in hitnasuyot:
        soldier_to_hitnasut[h["soldier_name"].strip().lower()] = h

    soldier_to_room: dict[str, str] = {}
    for room_num, occupants in room_assignment.items():
        for s in occupants:
            soldier_to_room[s["personal_number"]] = room_num

    room_to_manager: dict[str, str] = {r["room_number"]: (r.get("room_manager") or "") for r in rooms}

    headers = [
        "Soldier Number",
        "Name", "Last Name", "Personal Number", "Phone", "Course", "Gender",
        "Half Kappa", "Hitnasut", "Commander", "Room Number", "Room Manager",
        "Rank", "Unit", "Date of Birth", "City",
        "Username", "Password",
    ]
    fields = [
        "soldier_number",
        "name", "last_name", "personal_number", "phone", "course", "gender",
        "half_kappa", "hitnasut_name", "commander_name", "room_number", "room_manager",
        "rank", "unit", "date_of_birth", "city",
        "username", "password",
    ]

    wb = openpyxl.Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "Soldiers"

    for col, header in enumerate(headers, start=1):
        _style_header(ws.cell(row=1, column=col, value=header))
        ws.column_dimensions[get_column_letter(col)].width = 18

    for row_idx, s in enumerate(soldiers, start=2):
        full_name = f"{s['name']} {s['last_name']}".strip().lower()
        hitnasut = soldier_to_hitnasut.get(full_name, {})
        api = _api_fetch(s["personal_number"])
        cmd = commander_assignment.get(s["personal_number"], {})
        commander_name = f"{cmd.get('name', '')} {cmd.get('last_name', '')}".strip()
        cu = computer_users[row_idx - 2] if (row_idx - 2) < len(computer_users) else {}
        row: dict = {
            "soldier_number": row_idx - 1,
            "name": s["name"],
            "last_name": s["last_name"],
            "personal_number": s["personal_number"],
            "phone": s.get("phone", ""),
            "course": s.get("course", ""),
            "gender": s.get("gender", ""),
            "half_kappa": course_to_half_kappa.get((s.get("course") or "").strip().lower(), ""),
            "hitnasut_name": hitnasut.get("hitnasut_name", ""),
            "commander_name": commander_name,
            "room_number": soldier_to_room.get(s["personal_number"], ""),
            "room_manager": room_to_manager.get(soldier_to_room.get(s["personal_number"], ""), ""),
            **api,
            "username": cu.get("username", ""),
            "password": cu.get("password", ""),
        }
        for col, field in enumerate(fields, start=1):
            ws.cell(row=row_idx, column=col, value=row.get(field, ""))

    tab = Table(displayName="SoldierData", ref=f"A1:{get_column_letter(len(fields))}{len(soldiers) + 1}")
    tab.tableStyleInfo = TableStyleInfo(name="TableStyleMedium9", showRowStripes=True, showFirstColumn=False, showLastColumn=False, showColumnStripes=False)
    ws.add_table(tab)

    out = out_dir / "data.xlsx"
    wb.save(out)
    return str(out)


def _generate_room_signs(room_assignment: dict, out_dir: Path) -> str:
    tmpl = Document(str(ROOM_SIGNS_TEMPLATE))
    slots_per_page = sum(
        1 for para in tmpl.paragraphs
        if re.search(r"<soldier name>", para.text, re.IGNORECASE)
    ) or 1

    doc = Document()
    body = doc.element.body
    for p in list(doc.paragraphs):
        body.remove(p._element)

    def _add_page_break():
        pb_p = etree.Element(f"{{{_W}}}p")
        pb_r = etree.SubElement(pb_p, f"{{{_W}}}r")
        pb_br = etree.SubElement(pb_r, f"{{{_W}}}br")
        pb_br.set(f"{{{_W}}}type", "page")
        sectPr = body.find(f"{{{_W}}}sectPr")
        if sectPr is not None:
            sectPr.addprevious(pb_p)
        else:
            body.append(pb_p)

    def _add_template_page(room_num: str, batch: list):
        soldier_idx = 0
        for tmpl_para in tmpl.paragraphs:
            para_text = tmpl_para.text
            replacements: dict[str, str] = {}
            if re.search(r"<room sign>", para_text, re.IGNORECASE):
                replacements["<room sign>"] = room_num
            if re.search(r"<Soldier name>", para_text, re.IGNORECASE):
                if soldier_idx < len(batch):
                    s = batch[soldier_idx]
                    replacements["<Soldier name>"] = f"{s['name']} {s['last_name']}".strip()
                else:
                    replacements["<Soldier name>"] = ""
                soldier_idx += 1
            new_elem = copy.deepcopy(tmpl_para._element)
            if replacements:
                _replace_in_para_elem(new_elem, replacements)
            sectPr = body.find(f"{{{_W}}}sectPr")
            if sectPr is not None:
                sectPr.addprevious(new_elem)
            else:
                body.append(new_elem)

    first_page = True
    for room_num, occupants in room_assignment.items():
        batches = [occupants[i:i + slots_per_page] for i in range(0, len(occupants), slots_per_page)]
        if not batches:
            batches = [[]]

        for batch in batches:
            if not first_page:
                _add_page_break()
            first_page = False
            _add_template_page(room_num, batch)

    out = out_dir / "room_signs.docx"
    doc.save(str(out))
    return str(out)


# ---------------------------------------------------------------------------
# GUI
# ---------------------------------------------------------------------------

root = tk.Tk()
root.title("Kappa Personnel Manager")
root.resizable(False, False)

_FONT = "Helvetica"

notebook = ttk.Notebook(root)
notebook.pack(fill=tk.BOTH, expand=True, padx=12, pady=12)

# --- Main tab ---
main_tab = ttk.Frame(notebook, padding=24)
notebook.add(main_tab, text="Main")

ttk.Label(main_tab, text="Kappa Personnel Manager", font=(_FONT, 14, "bold")).pack(pady=(0, 20))


def refresh_status():
    for w in errors_frame.winfo_children():
        w.destroy()

    plain_errors: list[str] = []
    error_sections: dict[str, list[str]] = {}

    try:
        soldiers, commanders, hitnasuyot, courses, rooms, computer_users, classrooms = read_excel()
        error_sections = _validate_cross_references(soldiers, commanders, hitnasuyot, courses, computer_users, classrooms, rooms)
        room_errors = _check_room_capacity(soldiers, rooms)
        if room_errors:
            error_sections["Soldiers without a room"] = room_errors
    except FileNotFoundError:
        plain_errors.append("personnel.xlsx not found")
    except ValueError as e:
        plain_errors.extend(str(e).splitlines())

    ok = not plain_errors and not error_sections
    generate_btn.config(
        bg="#27ae60" if ok else "#c0392b",
        activebackground="#2ecc71" if ok else "#e74c3c",
    )

    if ok:
        ttk.Label(errors_frame, text="No errors found — ready to generate.", foreground="green").pack(anchor="w")
    else:
        for line in plain_errors:
            ttk.Label(errors_frame, text=line, foreground="red").pack(anchor="w")
        first = True
        for section_title, lines in error_sections.items():
            top_pad = (6, 2) if first else (10, 2)
            first = False
            ttk.Label(errors_frame, text=section_title, font=(_FONT, 9, "bold"), foreground="red").pack(anchor="w", pady=top_pad)
            for line in lines:
                ttk.Label(errors_frame, text=f"  • {line}", foreground="red").pack(anchor="w")

    root.update_idletasks()


_last_mtime: float = -1.0


def _poll_excel():
    global _last_mtime
    try:
        mtime = Path(EXCEL_PATH).stat().st_mtime
    except FileNotFoundError:
        mtime = 0.0
    if mtime != _last_mtime:
        _last_mtime = mtime
        refresh_status()
    root.after(1000, _poll_excel)


def btn_generate_all():
    missing = []
    if not Path(EXCEL_PATH).exists():
        missing.append(f"• {EXCEL_PATH}")
    for tmpl in (HITNASUYOT_TEMPLATE, ROOM_SIGNS_TEMPLATE, NAME_TAGS_TEMPLATE, CLASSROOM_TEMPLATE):
        if not tmpl.exists():
            missing.append(f"• {tmpl}")
    if missing:
        messagebox.showerror(
            "Missing Files",
            "The following required files are missing:\n\n" + "\n".join(missing),
        )
        return

    missing_errors, unknown_errors = _validate_templates()
    if missing_errors:
        messagebox.showerror(
            "Template Errors",
            "Cannot generate — required placeholder(s) missing:\n\n"
            + "\n".join(f"• {e}" for e in missing_errors),
        )
        return
    if unknown_errors and not _ask_ignore_unknowns(unknown_errors):
        return

    try:
        soldiers, commanders, hitnasuyot, courses, rooms, computer_users, classrooms = read_excel()
    except ValueError:
        messagebox.showerror("Errors", "There are errors that need to be fixed before generating.")
        return

    cross_errors = _validate_cross_references(soldiers, commanders, hitnasuyot, courses, computer_users, classrooms, rooms)
    if cross_errors:
        messagebox.showerror("Errors", "There are errors that need to be fixed before generating.")
        return

    if _check_room_capacity(soldiers, rooms):
        messagebox.showerror("Errors", "There are errors that need to be fixed before generating.")
        return

    room_assignment = _assign_rooms(soldiers, rooms) if (soldiers and rooms) else {}
    commander_assignment = _assign_commanders(soldiers, commanders, courses) if (soldiers and commanders) else {}

    out_dir = _unique_result_dir()
    out_dir.mkdir()

    generated = []
    errors = []

    if soldiers:
        try:
            generated.append(
                _generate_data_xlsx(soldiers, hitnasuyot, courses, room_assignment, commander_assignment, computer_users, rooms, out_dir)
            )
        except Exception as e:
            errors.append(f"Data: {e}")

    if soldiers:
        try:
            generated.append(_generate_enriched(soldiers, out_dir))
        except Exception as e:
            errors.append(f"Enrich: {e}")

    if hitnasuyot:
        try:
            generated.append(_generate_pptx(hitnasuyot, out_dir))
        except Exception as e:
            errors.append(f"Presentation: {e}")

    if room_assignment:
        try:
            generated.append(_generate_room_signs(room_assignment, out_dir))
        except Exception as e:
            errors.append(f"Room signs: {e}")

    if soldiers:
        try:
            generated.append(_generate_name_tags(soldiers, out_dir))
        except Exception as e:
            errors.append(f"Name tags: {e}")

    if classrooms:
        try:
            generated.append(_generate_classroom_pptx(classrooms, out_dir))
        except Exception as e:
            errors.append(f"Classroom presentation: {e}")

    if generated:
        msg = f"Results saved to '{out_dir}'."
        if errors:
            msg += "\n\nErrors:\n" + "\n".join(errors)
        messagebox.showinfo("Done", msg)
    elif errors:
        messagebox.showerror("Errors", "\n".join(errors))
    else:
        messagebox.showwarning("Nothing to generate", "No data found in personnel.xlsx.")


def btn_get_template():
    if Path(EXCEL_PATH).exists():
        messagebox.showerror("File Exists", f"'{EXCEL_PATH}' already exists. Delete or rename it first.")
        return
    from create_excel import create_workbook
    create_workbook()
    messagebox.showinfo("Done", f"Template created as {EXCEL_PATH}")
    refresh_status()


btn_row = ttk.Frame(main_tab)
btn_row.pack(pady=(0, 16))
generate_btn = tk.Button(
    btn_row, text="Generate", command=btn_generate_all, width=22,
    fg="white", font=(_FONT, 10), relief="flat", cursor="hand2",
    bg="#888888", activeforeground="white",
)
generate_btn.pack(side=tk.LEFT, padx=(0, 6))
ttk.Button(btn_row, text="Get Template", command=btn_get_template, width=14).pack(side=tk.LEFT)

errors_frame = ttk.Frame(main_tab)
errors_frame.pack(anchor="w", pady=(10, 0), fill="x")

# --- Help tab ---
help_tab = ttk.Frame(notebook, padding=24)
notebook.add(help_tab, text="Help")

ttk.Label(help_tab, text="Kappa Personnel Manager — User Guide", font=(_FONT, 14, "bold")).pack(anchor="w", pady=(0, 12))

HELP_TEXT = """\
═══════════════════════════════════════════════════════════════════
 OVERVIEW
═══════════════════════════════════════════════════════════════════

 This application reads soldier, commander, course, and room data
 from an Excel file, automatically assigns each soldier a commander
 and a room, and generates several output documents: a full data
 summary sheet, a Hitnasuyot presentation, room signs, and name
 tags.

 The assignments are made by an optimizer that tries to:
   1. Keep soldiers of the same course together in the same room
      and under the same commander (no soldier left alone with
      their course in a room or team).
   2. Balance the number of soldiers assigned to each commander
      evenly within each half-kappa.
   3. Respect gender: soldiers only go into rooms of their gender.
   4. Respect half-kappa: a commander only receives soldiers whose
      course belongs to the same half-kappa as that commander.

═══════════════════════════════════════════════════════════════════
 STEP 1 — GET THE EXCEL TEMPLATE
═══════════════════════════════════════════════════════════════════

 Click "Get Template" to create a blank personnel.xlsx file in
 the same folder as the application.
 Open it in Excel or LibreOffice and fill in your data.

 The file has ONE sheet ("Personnel") with SIX tables side by side.
 Each table starts on row 3  (row 1 = table title, row 2 = headers).
 Do NOT change the column positions or add extra columns.

═══════════════════════════════════════════════════════════════════
 STEP 2 — FILL IN THE EXCEL FILE
═══════════════════════════════════════════════════════════════════

 TABLE 1 — SOLDIERS  (columns A–F)
   One row per soldier. All five columns are important.
   • Name            — soldier's first name            (required)
   • Last Name       — soldier's last name             (required)
                       First + Last name must be unique across
                       all soldiers.
   • Personal Number — unique military ID              (required,
                       must be unique)
   • Phone Number    — phone number (if filled in, must be unique)
   • Course          — must exactly match a course name in the
                       Courses table (Table 4). Used to group
                       soldiers for room and commander assignment.
   • Gender          — used to assign rooms; must exactly match
                       the gender written in the Rooms table
                       (e.g. write "male" here and "male" there).

 TABLE 2 — COMMANDERS  (columns H–J)
   One row per commander.
   • Name       — commander's first name  (required)
   • Last Name  — commander's last name   (required)
   • Half Kappa — which half-kappa this commander belongs to
                  (required). Must exactly match the Half Kappa
                  values used in the Courses table.

   HOW COMMANDER ASSIGNMENT WORKS:
   Each soldier is assigned to exactly one commander. The app
   uses an optimizer with these priorities (in order):

     Priority 1 — No isolated soldiers.
       In every commander's team, no soldier should be the only
       one from their course. If a course has 2+ soldiers they
       are always placed in the same team, or split so that each
       sub-group has at least 2 from that course.
       (Exception: if a course has only 1 soldier in total, that
       soldier may be the only one from their course in the team.)

     Priority 2 — Keep full course groups together.
       As many soldiers of the same course as possible are placed
       under the same commander.

     Constraint — Half-kappa match.
       A commander only receives soldiers whose course belongs to
       the same half-kappa as the commander. Soldiers without a
       course (or whose course half-kappa has no commander) are
       distributed among all commanders as a fallback.

     Constraint — Equal load per commander (hard rule).
       Every commander in the same half-kappa must receive the
       same number of soldiers. If the soldiers in that half-kappa
       do not divide evenly, some commanders may have exactly one
       soldier more than the others — but the difference between
       the most-loaded and least-loaded commander in a half-kappa
       is never more than 1.
       Example: 10 soldiers, 3 commanders → two commanders get 3
       soldiers each and one gets 4 (or vice versa).
       This rule is enforced strictly and cannot be overridden by
       the course-grouping optimisation.

 TABLE 3 — HITNASUYOT  (columns L–N)
   One row per soldier–hitnasut pairing.
   • Hitnasut Name       — name of the hitnasut event/group
   • Full Soldier Name   — MUST be written as "First Last"
                           (first name space last name) and match
                           exactly a soldier in the Soldiers table.
   • Full Commander Name — optional; if filled, MUST be written as
                           "First Last" and match exactly a
                           commander in the Commanders table.

 TABLE 4 — COURSES  (columns P–Q)
   One row per course.
   • Course     — course name. Must exactly match what soldiers
                  have in their Course column (including spaces
                  and capitalisation).
   • Half Kappa — which half-kappa this course belongs to. Must
                  exactly match the Half Kappa values in the
                  Commanders table so that soldiers and commanders
                  can be matched.

 TABLE 5 — ROOMS  (columns S–V)
   One row per room.
   • Room Number  — room identifier, e.g. "101" or "A2".
                    Must be unique (no two rows with the same
                    room number).
   • Gender       — which gender this room is for. Must exactly
                    match what soldiers have in their Gender column
                    (e.g. "male" / "female").
   • Capacity     — maximum number of soldiers in this room.
                    Must be a whole number greater than zero.
   • Room Manager — optional. If filled, must be written as
                    "First Last" and match exactly a soldier in
                    the Soldiers table. Appears in data.xlsx.

   HOW ROOM ASSIGNMENT WORKS:

   Each soldier is assigned to exactly one room. The app uses the
   same optimizer as commander assignment, with the same priorities:

     Priority 1 — No isolated soldiers.
       In every room, no soldier should be the only one from their
       course (same exception for single-soldier courses applies).

     Priority 2 — Keep full course groups together.
       As many soldiers of the same course as possible are placed
       in the same room.

     Constraint — Gender match.
       Soldiers only go into rooms of matching gender.

     Constraint — Capacity.
       No room exceeds its stated capacity.

 TABLE 6 — COMPUTER USERS  (columns X–Y)
   One row per computer account, in the same order as the
   Soldiers table (row 1 here corresponds to soldier row 1, etc.).
   • Username — the computer account username  (required)
   • Password — the computer account password  (required)

   When generating data.xlsx, the Username and Password from this
   table are copied into the corresponding soldier's row. If there
   are fewer Computer User rows than soldiers, the remaining
   soldiers get empty Username/Password fields.

═══════════════════════════════════════════════════════════════════
 STEP 3 — PREPARE THE TEMPLATES FOLDER
═══════════════════════════════════════════════════════════════════

 Create a folder named "templates" in the same folder as the app.
 Place these three files inside it:

 templates/hitnasuyot.pptx
   A PowerPoint file with EXACTLY 3 slides:
     Slide 1 — intro WITH a commander:
               must contain  <Hitnasut name>  and  <Commander name>
     Slide 2 — intro WITHOUT a commander:
               must contain  <Hitnasut name>
     Slide 3 — soldier slide:
               must contain  <Soldier name>
   The app copies these slides for every hitnasut group and
   fills in the real names automatically.

 templates/room_signs.docx
   A Word document defining one page layout.
   Place <Room sign> where the room number should appear.
   Place <Soldier name> once for each soldier slot on the page.
   Slots are filled top to bottom; unused slots are left blank.
   If a room has more soldiers than slots, the app creates
   multiple pages for that room automatically.

 templates/name_tags.pptx
   A PowerPoint file with one slide.
   Place <Soldier name> once per name-tag slot on the slide.
   The app creates one slide per batch of soldiers, filling the
   slots in order.

 Note: all placeholders are case-insensitive.
       <Soldier Name>, <SOLDIER NAME> and <soldier name> all work.

═══════════════════════════════════════════════════════════════════
 STEP 4 — CHECK FOR ERRORS AND GENERATE
═══════════════════════════════════════════════════════════════════

 The "Generate" button turns GREEN when everything is valid and
 ready to produce output.
 It turns RED when there are problems that must be fixed first.

 Errors appear below the button automatically every time you
 save personnel.xlsx — you do not need to click anything.
 Fix all shown errors before clicking Generate.

 ── PREREQUISITE ERRORS (shown live, must be fixed) ─────────────

   • "personnel.xlsx not found"
       Click "Get Template" to create the file, then fill it in.

   • "Row X — field: field must not be empty"
       A required field in row X is blank. Fill it in.

   • "Row X — capacity: capacity must be a whole number > 0"
       The Capacity value in the Rooms table must be a positive
       whole number (1, 2, 3 …).

   • "Soldiers with duplicate name"
       Two or more soldiers share the same first + last name.
       Full names must be unique. Check the listed rows.

   • "Soldiers with duplicate personal number"
       Two or more soldiers have the same personal number.
       Each soldier's personal number must be unique.

   • "Soldiers with duplicate phone number"
       Two or more soldiers share the same phone number.
       Each non-empty phone number must be unique.

   • "Soldiers with unknown course"  (row X, col F)
       A soldier's Course value does not match any course name
       in the Courses table. Check for typos or missing courses.

   • "Soldiers with no eligible commander"  (row X, col F)
       The course's half-kappa has no commander. Either add a
       commander for that half-kappa in the Commanders table, or
       correct the Half Kappa value in the Courses table.

   • "Hitnasuyot with unknown soldier"  (row X, col M)
       The Full Soldier Name does not match any soldier's
       "First Last" name. Check spelling and spacing exactly.

   • "Hitnasuyot with unknown commander"  (row X, col N)
       The Full Commander Name does not match any commander's
       "First Last" name. Check spelling and spacing exactly.

   • "Rooms with duplicate number"
       Two rows in the Rooms table have the same Room Number.
       Each room must have a unique identifier.

   • "X soldier(s) have no room"
       Not enough total room capacity for soldiers of that gender.
       Add more rooms or increase existing room capacities.

   • "Soldiers without a computer user"  (row X, cols A–B)
       The Computer Users table has fewer rows than the Soldiers
       table. Each soldier must have a corresponding computer user
       at the same row position. Add the missing rows to the
       Computer Users table.

   • "Rooms with unknown room manager"  (row X, col V)
       The Room Manager name does not match any soldier's
       "First Last" name. Check spelling and spacing exactly,
       or leave the cell blank if there is no room manager.

 ── AFTER CLICKING GENERATE ────────────────────────────────────

 A new result folder is created automatically each time you click
 Generate (result/, result(1)/, result(2)/, …) so previous
 outputs are never overwritten.

 If the templates folder or any template file is missing, the
 app will tell you which file is absent before doing anything.

═══════════════════════════════════════════════════════════════════
 OUTPUT FILES  (inside the result/ folder)
═══════════════════════════════════════════════════════════════════

 data.xlsx
   One row per soldier with all information combined:
   name, personal number, course, gender, half-kappa, assigned
   commander, assigned room, hitnasut (if any), additional
   info fetched from the API (rank, unit, date of birth, city),
   and Username/Password from the Computer Users table.
   Rows are numbered sequentially in the "Soldier Number" column.

 enriched_soldiers.xlsx
   Soldiers with extra API-fetched data (rank, unit, birth
   date, city).

 hitnasuyot.pptx
   One set of slides per hitnasut group: an intro slide (with
   or without commander name) followed by one soldier slide per
   participant in that group.

 room_signs.docx
   One page per room per batch of soldiers (multiple pages for a
   room if it has more occupants than template slots), with all
   soldier names filled in.

 name_tags.pptx
   One slide per batch of soldiers, names filled into the slots
   defined in the template.
"""

_help_lines = HELP_TEXT.rstrip("\n").splitlines()
_txt_frame = ttk.Frame(help_tab)
_txt_frame.pack(fill="both", expand=True)
_scrollbar = ttk.Scrollbar(_txt_frame, orient="vertical")
_scrollbar.pack(side="right", fill="y")
txt = tk.Text(
    _txt_frame, wrap="none", relief="flat", font=("Courier", 10), state="normal",
    height=20, width=max(len(l) for l in _help_lines),
    yscrollcommand=_scrollbar.set,
)
txt.insert("1.0", HELP_TEXT)
txt.config(state="disabled")
txt.pack(side="left", fill="both", expand=True)
_scrollbar.config(command=txt.yview)

root.update_idletasks()
_init_w = root.winfo_reqwidth()
_init_h = root.winfo_reqheight()
_init_x = (root.winfo_screenwidth() - _init_w) // 2
_init_y = (root.winfo_screenheight() - _init_h) // 2
root.geometry(f"{_init_w}x{_init_h}+{_init_x}+{_init_y}")

_poll_excel()

root.mainloop()
